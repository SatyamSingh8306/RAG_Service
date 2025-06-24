import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import os
from typing import List, Dict, Any, Optional
from pathlib import Path

# Imports for semantic chunking
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.document_loaders import CSVLoader, Docx2txtLoader, WebBaseLoader
from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
import pandas as pd
import json
import zipfile
import xml.etree.ElementTree as ET
from docx import Document
import requests
import re

# Local application imports
from backend.app.core.config import settings
from backend.app.services.vstore_svc import VectorStoreService

class DocParserFastService:
    """
    Service to parse documents (PDFs, images), extract text,
    perform STRUCTURE-AWARE semantic chunking, and add chunks to the vector store.
    """
    def __init__(self, vector_store_service: VectorStoreService):
        self.vector_store_service = vector_store_service
        self.ocr_processor = None  # Will be initialized on first use

        # Use Instructor embeddings model for more citation-aware and instruction-aligned chunks
        # light_model = sentence-transformers/all-MiniLM-L6-v2
        # powerful and free model = hkunlp/instructor-xl
        try:
            self.embeddings = HuggingFaceEmbeddings(
                model_name="sentence-transformers/all-MiniLM-L6-v2",  # Free and powerful model from HuggingFace
                model_kwargs={'device': 'cpu'},
                encode_kwargs={"normalize_embeddings": True}
            )

            # NOTE: Using semantic chunking with INSTRUCTOR may require a custom splitter; fallback used here
            self.semantic_splitter = RecursiveCharacterTextSplitter(
                chunk_size=800,  # smaller for citation-aware precision
                chunk_overlap=200,
                length_function=len,
                is_separator_regex=False,
                separators=["\n\n", "\n", ". ", "? ", "! ", ",", " ", ""]
            )
        except Exception as e:
            print(f"Error initializing embedding model: {e}. Using fallback splitter.")
            self.embeddings = None
            self.semantic_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len,
                is_separator_regex=False,
                separators=["\n\n", "\n", ". ", "? ", "! ", ",", " ", ""]
            )

    def _is_heading(self, text: str) -> bool:
        """
        Detect if a line/paragraph is a heading based on heuristic:
        - Mostly uppercase or capitalized and short (< 8 words)
        """
        words = text.strip().split()
        if len(words) <= 8 and sum(1 for w in words if w.isupper()) >= len(words) * 0.6:
            return True
        return False

    def _extract_sections_from_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from PDF and groups paragraphs into sections by detecting headings.
        Returns list of sections with 'title', 'page_number', and 'paragraphs'.
        """
        sections = []
        doc = fitz.open(file_path)
        section_counter = 0
        for page_index in range(len(doc)):
            page = doc.load_page(page_index)
            blocks = page.get_text("blocks", sort=True)
            current_section = {"title": f"page_{page_index+1}_untitled", "page_number": page_index+1, "paragraphs": []}
            for block in blocks:
                if block[6] != 0:
                    continue
                text = block[4].replace('\r', ' ').replace('\n', ' ').strip()
                if not text:
                    continue
                if self._is_heading(text):
                    # start new section
                    if current_section["paragraphs"]:
                        sections.append(current_section)
                    section_counter += 1
                    current_section = {"title": text, "page_number": page_index+1, "paragraphs": [], "section_index": section_counter}
                else:
                    current_section["paragraphs"].append(text)
            if current_section["paragraphs"]:
                # ensure section_index for untitled sections
                if "section_index" not in current_section:
                    section_counter += 1
                    current_section["section_index"] = section_counter
                sections.append(current_section)
        doc.close()
        return sections

    def _extract_text_from_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from PPTX and groups paragraphs into sections by detecting headings.
        Returns list of sections with 'title', 'page_number', and 'paragraphs'.
        """
        prs = Presentation(file_path)
        sections = []
        section_counter = 0

        for i, slide in enumerate(prs.slides):
            current_section = {
                "title": f"slide_{i+1}_untitled",
                "page_number": i + 1,
                "paragraphs": []
            }

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip().replace('\r', ' ').replace('\n', ' ')
                    if not text:
                        continue
                    if self._is_heading(text):
                        # Start new section
                        if current_section["paragraphs"]:
                            sections.append(current_section)
                        section_counter += 1
                        current_section = {
                            "title": text,
                            "page_number": i + 1,
                            "paragraphs": [],
                            "section_index": section_counter
                        }
                    else:
                        current_section["paragraphs"].append(text)

                elif isinstance(shape, GraphicFrame) and shape.has_chart:
                    chart = shape.chart  # safe to access now
                    for plot in chart.plots:
                        categories = [c.label for c in plot.categories]
                        for series in plot.series:
                            values = tuple(series.values)
                            chart_text = (
                                f"Chart ({chart.chart_type}): Series '{series.name}', "
                                f"categories {categories}, values {values}."
                            )
                            current_section["paragraphs"].append(chart_text)

            if current_section["paragraphs"]:
                if "section_index" not in current_section:
                    section_counter += 1
                    current_section["section_index"] = section_counter
                sections.append(current_section)

        return sections


    def _extract_text_from_csv(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from a CSV file and groups rows into sections.
        Returns list of sections with 'title', 'page_number', and 'paragraphs'.
        """
        sections = []
        section_counter = 0

        loader = CSVLoader(file_path=file_path)
        # documents = loader.load_and_split(text_splitter=self.semantic_splitter)
        documents = loader.load()

        for i, doc in enumerate(documents):
            text = doc.page_content.strip()
            if not text:
                continue
            current_section = {
                "title": f"page_{i+1}_untitled",
                "page_number": i + 1,
                "paragraphs": [text]
            }
            section_counter += 1
            current_section["section_index"] = section_counter
            sections.append(current_section)

        return sections

    def _extract_text_from_excel(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from Excel file and groups each sheet as a section.
        Returns list of sections with 'title', 'page_number', and 'paragraphs'.
        """
        sections = []
        section_counter = 0

        excel_file = pd.ExcelFile(file_path)

        for i, sheet_name in enumerate(excel_file.sheet_names):
            df = pd.read_excel(file_path, sheet_name=sheet_name)

            # Convert DataFrame to text (one big paragraph)
            text_content = df.to_string(index=False).strip()
            if not text_content:
                continue

            current_section = {
                "title": sheet_name if sheet_name else f"sheet_{i+1}_untitled",
                "page_number": i + 1,
                "paragraphs": [text_content],
            }

            section_counter += 1
            current_section["section_index"] = section_counter
            sections.append(current_section)

        return sections
    
    def _extract_text_from_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from DOCX file and split into sections.
        Returns list of sections with 'title', 'page_number', and 'paragraphs'.
        """
        sections = []
        section_counter = 0
        loader = Docx2txtLoader(file_path=file_path)
        documents = loader.load()
        
        for i, doc in enumerate(documents):
            text = doc.page_content.strip()
            if not text:
                continue
            section_counter += 1
            sections.append({
                "title": f"page_{i+1}_untitled",
                "page_number": i + 1,
                "paragraphs": [text],
                "section_index": section_counter
            })

        # with zipfile.ZipFile(file_path, 'r') as docx:
        #     chart_files = [f for f in docx.namelist() if f.startswith('word/charts/chart')]
        #     for chart_file in chart_files:
        #         chart_xml = docx.read(chart_file)
        #         root = ET.fromstring(chart_xml)
        #         namespaces = {
        #             'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        #         }
        #         chart_data = []
        #         for pt in root.findall('.//c:pt', namespaces):
        #             idx = pt.get('idx')
        #             value = pt.find('c:v', namespaces).text
        #             chart_data.append((idx, value))
                
        #         # Format chart data into a readable paragraph
        #         chart_para = f"Chart '{chart_file}': " + ', '.join([f"Index {idx}: {value}" for idx, value in chart_data])
                
        #         # Append the chart data to the paragraphs of the current section
        #         if sections:
        #             sections[-1]["paragraphs"].append(chart_para)
        return sections

    def _extract_text_from_json(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Handles JSON files containing multiple JSON objects (one per line or concatenated).
        Extracts structured sections with 'title', 'page_number', 'paragraphs', and 'section_index'.
        """
        sections = []
        section_counter = 0

        with open(file_path, "r", encoding="utf-8") as f:
            for i, line in enumerate(f):
                line = line.strip()
                if not line:
                    continue
                try:
                    item = json.loads(line)
                except json.JSONDecodeError as e:
                    print(f"Skipping line {i+1} due to JSON decode error: {e}")
                    continue

                title = item.get("title") or f"section_{i+1}_untitled"
                text_parts = []

                for key, value in item.items():
                    if isinstance(value, str):
                        text_parts.append(value)
                    elif isinstance(value, list):
                        text_parts.extend([str(v) for v in value if isinstance(v, str)])

                paragraph_text = " ".join(text_parts).strip()
                if not paragraph_text:
                    continue

                section_counter += 1
                current_section = {
                    "title": title,
                    "page_number": section_counter,
                    "paragraphs": [paragraph_text],
                    "section_index": section_counter
                }
                sections.append(current_section)

        return sections


    def _extract_text_from_txt(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from TXT file and splits by headings or paragraphs.
        Returns list of sections with 'title', 'page_number', and 'paragraphs'.
        """
        sections = []
        section_counter = 0

        with open(file_path, "r", encoding="utf-8") as f:
            lines = [line.strip() for line in f if line.strip()]

        current_section = {
            "title": "page_1_untitled",
            "page_number": 1,
            "paragraphs": []
        }

        for line in lines:
            if self._is_heading(line):
                if current_section["paragraphs"]:
                    sections.append(current_section)
                section_counter += 1
                current_section = {
                    "title": line,
                    "page_number": 1,
                    "paragraphs": [],
                    "section_index": section_counter
                }
            else:
                current_section["paragraphs"].append(line)

        if current_section["paragraphs"]:
            if "section_index" not in current_section:
                section_counter += 1
                current_section["section_index"] = section_counter
            sections.append(current_section)

        return sections
    
    def _extract_text_from_html(self, url: str) -> List[Dict[str, Any]]:
        sections =[]
        section_counter = 0
        try:
            loader = WebBaseLoader([url]) # WebBaseLoader expects a list of URLs
            documents = loader.load()
            for i,doc in enumerate(documents):
                text = doc.page_content.strip()
                if text:
                    title = doc.metadata.get("title", f"URL Content {i+1}")
                    current_section = {"title": title, "page_number": 1, "paragraphs": [text], "section_index": section_counter + 1}
                    sections.append(current_section)
                    section_counter += 1
        except Exception as e:
            print(f"Error in loading the html file {url}: {e}")
        return sections
    
    def extract_doc_id_from_url(self, url: str) -> str:
        """
        Extracts the Google Docs document ID from a full URL.
        """
        match = re.search(r"/d/([a-zA-Z0-9-_]+)", url)
        if match:
            return match.group(1)
        else:
            raise ValueError("Invalid Google Docs URL format.")


    def _extract_text_from_gdoc_url(self, url: str) -> List[Dict[str, Any]]:
        """
        Extracts text from a public Google Docs URL (exported as plain text).
        Returns list of sections with 'title', 'page_number', and 'paragraphs'.
        """
        sections = []
        section_counter = 0

        doc_id = self.extract_doc_id_from_url(url)
        export_url = f"https://docs.google.com/document/d/{doc_id}/export?format=txt"
        response = requests.get(export_url)

        if response.status_code != 200:
            raise Exception(f"Failed to fetch Google Doc content: {response.status_code}")

        lines = [line.strip() for line in response.text.splitlines() if line.strip()]

        current_section = {
            "title": "page_1_untitled",
            "page_number": 1,
            "paragraphs": []
        }

        for line in lines:
            if self._is_heading(line):
                if current_section["paragraphs"]:
                    sections.append(current_section)
                section_counter += 1
                current_section = {
                    "title": line,
                    "page_number": 1,
                    "paragraphs": [],
                    "section_index": section_counter
                }
            else:
                current_section["paragraphs"].append(line)

        if current_section["paragraphs"]:
            if "section_index" not in current_section:
                section_counter += 1
                current_section["section_index"] = section_counter
            sections.append(current_section)

        return sections
    
    def _extract_text_from_gsheet_url(self, url: str) -> List[Dict[str, Any]]:
        """
        Extracting the data from google sheet (using url importing it as csv)
        RETURN a list of dict which contain data , id and more.
        """
        sections = []
        section_counter = 0

        # Convert the Google Sheets URL to CSV export format
        if "/edit#gid=" in url:
            base_url, gid_part = url.split("/edit#gid=")
            sheet_id = base_url.split("/d/")[1]
            gid = gid_part
        else:
            raise ValueError("Invalid Google Sheets URL format.")

        export_url = f"https://docs.google.com/spreadsheets/d/{sheet_id}/export?format=csv&gid={gid}"
        response = requests.get(export_url)

        if response.status_code != 200:
            raise Exception(f"Failed to fetch Google Sheet content: {response.status_code}")

        # Load into pandas DataFrame
        df = pd.read_csv(pd.compat.StringIO(response.text))

        current_section = {
            "title": "sheet_data",
            "page_number": 1,
            "paragraphs": []
        }

        for _, row in df.iterrows():
            # Create a readable paragraph from the row
            paragraph = "; ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val) and str(val).strip()])
            if paragraph:
                current_section["paragraphs"].append(paragraph)

        if current_section["paragraphs"]:
            section_counter += 1
            current_section["section_index"] = section_counter
            sections.append(current_section)

        return sections


    def _extract_text_from_image(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from an image using OCR.
        Treats the whole image as one section.
        """
        sections = []
        try:
            img = Image.open(file_path)
            text_from_ocr = pytesseract.image_to_string(img)
            
            if text_from_ocr.strip():
                # Split text into paragraphs
                paragraphs = [p.strip() for p in text_from_ocr.split('\n\n') if p.strip()]
                if not paragraphs:  # If no double newlines, try single newlines
                    paragraphs = [p.strip() for p in text_from_ocr.split('\n') if p.strip()]
                
                if paragraphs:
                    sections.append({
                        "title": "ocr_image_section",
                        "page_number": 1,
                        "paragraphs": paragraphs,
                        "section_index": 1
                    })
            print(f"Extracted text using OCR from image: {os.path.basename(file_path)}")
        except pytesseract.TesseractNotFoundError:
            print("TesseractNotFoundError: Tesseract is not installed or not in your PATH.")
            print("Please install Tesseract OCR and ensure it's accessible.")
        except Exception as e:
            print(f"Error extracting text from image {os.path.basename(file_path)} using OCR: {e}")
        return sections

    def _perform_semantic_chunking(self, text: str) -> List[str]:
        """
        Perform semantic chunking on text using embeddings, fallback if needed.
        Here we use RecursiveCharacterTextSplitter for INSTRUCTOR-compatible chunking.
        """
        try:
            return self.semantic_splitter.split_text(text)
        except Exception as e:
            print(f"Chunking failed: {e}")
            return [text]  # fallback to whole text if even fallback splitter fails

    def process_document(self, file_path: str, source_doc_id: str, collection_name: Optional[str] = None) -> bool:
        """Process a document using fast rule-based chunking."""
        try:
            # Extract text based on file type
            file_extension = Path(file_path).suffix.lower()
            if file_extension == ".pdf":
                sections = self._extract_sections_from_pdf(file_path)
            elif file_extension == ".pptx":
                sections = self._extract_text_from_pptx(file_path)
            elif file_extension == ".csv":
                sections = self._extract_text_from_csv(file_path)
            elif file_extension in [".xlsx", ".xls"]:
                sections = self._extract_text_from_excel(file_path)
            elif file_extension == ".docx":
                sections = self._extract_text_from_docx(file_path)
            elif file_extension == ".json":
                sections = self._extract_text_from_json(file_path)
            elif file_extension == ".txt":
                sections = self._extract_text_from_txt(file_path)
            elif file_extension in [".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif"]:
                sections = self._extract_text_from_image(file_path)
            else:
                print(f"Unsupported file type: {file_extension}")
                return False

            if not sections:
                print(f"No text extracted from {file_path}")
                return False

            # Process the text
            all_chunks, all_metadatas, all_ids = [], [], []
            global_chunk_counter = 1
            for sec in sections:
                section_text = "\n\n".join(sec["paragraphs"])
                chunks = self._perform_semantic_chunking(section_text)
                for idx, chunk in enumerate(chunks):
                    sec_idx = sec.get("section_index", sec["page_number"])
                    chunk_id = f"{source_doc_id}_sec{sec_idx}_chunk{global_chunk_counter}"
                    metadata = {
                        "source_doc_id": source_doc_id,
                        "file_name": Path(file_path).name,
                        "section_title": sec["title"],
                        "page_number": sec["page_number"],
                        "section_index": sec_idx,
                        "chunk_index": global_chunk_counter,
                        "paragraph_number_in_page": idx + 1
                    }
                    all_chunks.append(chunk)
                    all_metadatas.append(metadata)
                    all_ids.append(chunk_id)
                    global_chunk_counter += 1

            # Add documents to vector store
            success = self.vector_store_service.add_documents(
                chunks=all_chunks,
                metadatas=all_metadatas,
                doc_ids=all_ids,
                collection_name=collection_name
            )

            return success
        except Exception as e:
            print(f"Error processing document {file_path}: {e}")
            return False

    def process_url(self, url: str, source_doc_id: str, collection_name: Optional[str] = None) -> bool:
        """Process a document from a URL using fast rule-based chunking."""
        try:
            # Check if the URL is a Google Doc and route accordingly
            if "docs.google.com/document/" in url:
                print(f"Detected Google Doc URL: {url}")
                sections = self._extract_text_from_gdoc_url(url)
            elif "docs.google.com/spredsheet/" in url:
                sections = self._extract_text_from_gsheet_url(url)
            else:
                print(f"Detected standard web URL: {url}")
                sections = self._extract_text_from_html(url)

            if not sections:
                print(f"No text extracted from {url}")
                return False

            # Process the text
            all_chunks, all_metadatas, all_ids = [], [], []
            global_chunk_counter = 1
            for sec in sections:
                section_text = "\n\n".join(sec["paragraphs"])
                chunks = self._perform_semantic_chunking(section_text)
                for idx, chunk in enumerate(chunks):
                    sec_idx = sec.get("section_index", sec["page_number"])
                    chunk_id = f"{source_doc_id}_sec{sec_idx}_chunk{global_chunk_counter}"
                    metadata = {
                        "source_doc_id": source_doc_id,
                        "file_name": url,  # Use the URL as the identifier
                        "section_title": sec["title"],
                        "page_number": sec["page_number"],
                        "section_index": sec_idx,
                        "chunk_index": global_chunk_counter,
                        "paragraph_number_in_page": idx + 1
                    }
                    all_chunks.append(chunk)
                    all_metadatas.append(metadata)
                    all_ids.append(chunk_id)
                    global_chunk_counter += 1

            # Add documents to vector store
            success = self.vector_store_service.add_documents(
                chunks=all_chunks,
                metadatas=all_metadatas,
                doc_ids=all_ids,
                collection_name=collection_name
            )

            return success
        except Exception as e:
            print(f"Error processing URL {url}: {e}")
            return False

def process_document_background(
    file_path: str, 
    source_doc_id: str, 
    doc_parser_svc_instance: DocParserFastService,
    serial_no: Optional[int] = None, 
    total_count: Optional[int] = None,
    collection_name: Optional[str] = None
):
    parser_name = doc_parser_svc_instance.__class__.__name__
    progress_log = f"Document {serial_no}/{total_count}" if serial_no and total_count else "Single document"
    
    print(f"Background task started for: {Path(file_path).name}, Source ID: {source_doc_id}, Parser: {parser_name}, ({progress_log})")
    try:
        success = doc_parser_svc_instance.process_document(
            file_path=file_path, 
            source_doc_id=source_doc_id,
            collection_name=collection_name
        )
        if success:
            print(f"Background processing completed successfully for {source_doc_id} ({Path(file_path).name}) using {parser_name}")
        else:
            print(f"Background processing (using {parser_name}) had issues or no chunks generated for {source_doc_id} ({Path(file_path).name})")
    except Exception as e:
        print(f"Error during background document processing for {source_doc_id} ({Path(file_path).name}) (using {parser_name}): {e}")


def process_url_background(
    url: str,
    source_doc_id: str,
    doc_parser_svc_instance: DocParserFastService,
    serial_no: Optional[int] = None,
    total_count: Optional[int] = None,
    collection_name: Optional[str] = None
):
    parser_name = doc_parser_svc_instance.__class__.__name__
    progress_log = f"URL {serial_no}/{total_count}" if serial_no and total_count else "Single URL"

    print(f"Background task started for: {url}, Source ID: {source_doc_id}, Parser: {parser_name}, ({progress_log})")
    try:
        success = doc_parser_svc_instance.process_url(
            url=url,
            source_doc_id=source_doc_id,
            collection_name=collection_name
        )
        if success:
            print(f"Background processing completed successfully for {source_doc_id} ({url}) using {parser_name}")
        else:
            print(f"Background processing (using {parser_name}) had issues or no chunks generated for {source_doc_id} ({url})")
    except Exception as e:
        print(f"Error during background document processing for {source_doc_id} ({url}) (using {parser_name}): {e}")


if __name__ == "__main__":
    process_document_background(file_path=r"C:\Users\hp\Downloads\Resume (1).docx")